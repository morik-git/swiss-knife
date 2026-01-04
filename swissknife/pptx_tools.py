from __future__ import annotations

from pathlib import Path
from typing import Iterable

from pptx import Presentation


def merge_pptx(inputs: Iterable[str], output: str) -> Path:
    paths = [Path(p) for p in inputs]
    if not paths:
        raise ValueError("inputs must not be empty")

    base = Presentation(paths[0])
    for p in paths[1:]:
        prs = Presentation(p)
        for slide in prs.slides:
            blank_layout = base.slide_layouts[6]
            new_slide = base.slides.add_slide(blank_layout)
            for shape in slide.shapes:
                el = shape.element
                new_slide.shapes._spTree.insert_element_before(el, "p:extLst")

    out_path = Path(output)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    base.save(out_path)
    return out_path
